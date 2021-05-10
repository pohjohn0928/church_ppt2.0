import os
from pptx import Presentation
from pptx.util import Inches, Pt,Cm
from pptx.dml.color import RGBColor
import requests
from bs4 import BeautifulSoup
import time
import threading

class ReadPdfFile:
    def __init__(self):
        self.root = os.path.dirname(__file__)

    def getChieseScripture(self,scriptures):
        if scriptures == []:
            return []
        else:
            dic = BibleApi().english_chinese_dic
            return_scripture = []
            for scripture in scriptures:
                i = scripture.split(' ')[0]
                i = dic[i]
                j = scripture.split(' ')[1]
                return_scripture.append(i + ' ' + j)
            return return_scripture

    def getClosingSong(self,closingSongName):
        path = os.path.join(self.root,f'closingSong/{closingSongName}.txt')
        try:
            file = open(path,encoding='utf-8')
            lyrics = ""
            for f in file:
                lyrics += f
            closingSong_chinese = lyrics.split('(English)')[0].split('\n')
            closingSong_english = lyrics.split('(English)')[1].split('\n')
            del closingSong_chinese[-1]
            del closingSong_english[0]
            closingSong = {"closingSong_chinese":closingSong_chinese,"closingSong_english" : closingSong_english}
            return closingSong
        except:
            print("Can't find out the closingSongName !")

    def getBlessingSong(self):
        path = self.root + '/blessingSong/blessing_song.txt'
        file = open(path)
        blessing_song = []
        for lyrics in file:
            blessing_song.append(lyrics.replace('\n',''))
        return blessing_song

class BibleApi:
    def __init__(self):
        root = os.path.dirname(__file__)

        self.englishToChinese = root + '/Books/englishToChinese.txt'
        self.bibleVersion = root + "/bibleVersion/bibleVersion.txt"
        self.englishBook = root + "/Books/englishBook.txt"
        self.chineseBook = root + "/Books/chineseBook.txt"
        self.setBibleVersion()
        self.setEnglishBook()
        self.setChineseBook()
        self.english_to_chinese()
        self.englishVerseBaseUrl = "https://biblia-api-pdf.herokuapp.com/api?method=bible&"
        self.chineseVerseBaseUrl = "http://ibibles.net/quote.php?cut-"

    def english_to_chinese(self):
        self.english_chinese_dic = {}
        file = open(self.englishToChinese,encoding='utf-8')
        for scripture in file:
            self.english_chinese_dic[scripture.split(' ')[0]] = scripture.split(' ')[1].replace('\n','')


    def getEnglishVerse(self,bibleVersion,book,chapter,startVerse,endVerse):
        startVerse = int(startVerse)
        endVerse = int(endVerse)
        return_verses = []
        book = book.replace(' ','-')
        url = f"https://www.biblestudytools.com/{bibleVersion}/{book}/{chapter}.html"
        response = requests.get(url)
        soup = BeautifulSoup(response.text, "html.parser")

        for i in range(startVerse,endVerse + 1):
            result = soup.find("span",class_=f"verse-{i}")
            try:
                while(1):
                    result.a.decompose()
            except:
                pass
            return_verses.append(f"<sv>{i} {result.text.strip()}")
        return return_verses


    def getChineseVerse(self,book,chapter,startVerse,endVerse):
        url = self.chineseVerseBaseUrl + f"{book}/{chapter}:{startVerse}-{endVerse}"
        r = requests.get(url)
        soup = BeautifulSoup(r.text, 'html.parser')
        return_verses = soup.find('body').get_text().split('\n')
        return_verses = ["<sv>" + verse for verse in return_verses if verse != ""]
        # return_verses = [" ".join(verse.split(" ")[1:]) for verse in return_verses if verse != ""]

        return return_verses


    def setBibleVersion(self):
        self.bibleVersionDic = {}
        file = open(self.bibleVersion)
        for bibleVersion in file:
            self.bibleVersionDic[bibleVersion.replace("\n","").split(",")[0]] = bibleVersion.replace("\n","").split(",")[1]

    def setEnglishBook(self):
        self.englishBookDic = {}
        file = open(self.englishBook)
        for book in file:
            self.englishBookDic[book.replace("\n","").split(" ",2)[2]] = book.replace("\n","").split(" ",2)[1]


    def setChineseBook(self):
        self.chineseBookDic = {}
        file = open(self.chineseBook,"r",encoding="utf-8")
        for book in file:
            self.chineseBookDic[book.replace("\n","").split(" ")[1]] = book.replace("\n","").split(" ")[0]


class MakePPT(threading.Thread):
    def __init__(self, data):
        threading.Thread.__init__(self)
        self.prs = Presentation()
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)
        self.data = data
        self.word_limit = 69
        self.word_limit_chinese = 155
        self.layout = self.prs.slide_layouts[6]
        self.start_verse_token = "<sv>"

    def run(self):
        self.insertScriptureData()

    def insertScriptureData(self):
        start = time.time()
        englishScrpitureReading = self.data["englishScrpitureReading"]
        chineseScrpitureReading = self.data["chineseScrpitureReading"]
        englishScrpitureInSermon = self.data["englishScrpitureInSermon"]
        chineseScrpitureInSermon = self.data["chineseScrpitureInSermon"]
        self.addCalvaryImg()
        self.addAnnocement()
        self.representAtPPTScrpitureReading(englishScrpitureReading)
        self.representAtPPTScrpitureReading(chineseScrpitureReading,'chinese')
        self.addCalvaryImg()
        self.representAtPPTScrpitureInSermon(englishScrpitureInSermon,chineseScrpitureInSermon)
        self.addCalvaryImg()
        self.addClosingSong()
        self.addCalvaryImg()
        self.addBlessingSong()
        self.addCalvaryImg()
        self.prs.save(f'churchPPT{self.data["date"]}.pptx')
        end = time.time()
        print("Make PPT cost：%f sec" % (end - start))


    def addAnnocement(self):
        if self.data['annocement'] == []:
            pass

        else:
            for annoce in self.data['annocement']:
                img_path = os.path.dirname(__file__) + f'/annocement/{annoce}.png'
                slide = self.prs.slides.add_slide(self.layout)
                top = Inches(0)
                left = Inches(0)
                height = Inches(9)
                slide.shapes.add_picture(img_path, left, top, height=height)
            self.addCalvaryImg()

    def representAtPPTScrpitureReading(self,info,language = 'english'):
        if language == 'english':
            ScrpitureInSermon_title_slide_layout = self.prs.slide_layouts[0]
            slide = self.prs.slides.add_slide(ScrpitureInSermon_title_slide_layout)
            title = slide.shapes.title

            title.width = Cm(40)
            title.top = Inches(4)
            title.text = "Scrpiture Reading"
            title.text_frame.paragraphs[0].font.size = Pt(70)
            for index in range(len(info["verses"])):
                scripture = info["verses"][index]
                bibleVersion = info["bibleVersion"][index]
                verses = self.getEnglishBibleVerses(scripture,bibleVersion)
                if bibleVersion != 'NKJV' and bibleVersion != 'nkjv':
                    scripture += ' ' + f"({bibleVersion})"
                self.addPage(scripture,verses)
        else:
            for scripture in info:
                verses = self.getChineseBibleVerses(scripture)
                self.addPage(scripture,verses,'chinese')

    def representAtPPTScrpitureInSermon(self,englishScrpitureInSermon,chineseScrpitureInSermon):
        englishVerses = []
        chineseVerses = []
        for index in range(len(englishScrpitureInSermon["verses"])):
            scripture = englishScrpitureInSermon["verses"][index]
            bibleVersion = englishScrpitureInSermon["bibleVersion"][index]
            verses = self.getEnglishBibleVerses(scripture,bibleVersion)
            if verses == None:
                verses = ["None"]
            englishVerses.append(verses)
        for chineseScrpiture in chineseScrpitureInSermon:
            verses = self.getChineseBibleVerses(chineseScrpiture)
            chineseVerses.append(verses)

        ScrpitureInSermon_title_slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(ScrpitureInSermon_title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Sermon Title"
        title.width = Cm(40)
        title.top = Inches(4)
        subtitle.text = self.data["sermonTitle"]
        subtitle.width = Cm(40)
        subtitle.top = Inches(5)
        title.text_frame.paragraphs[0].font.size = Pt(70)
        subtitle.text_frame.paragraphs[0].font.size = Pt(50)
        for i in range(len(englishScrpitureInSermon["verses"])):
            slide = self.prs.slides.add_slide(self.layout)
            tf = self.initTxBox(slide)
            chapter = englishScrpitureInSermon["verses"][i]
            if englishScrpitureInSermon['bibleVersion'][i] != 'NKJV' and englishScrpitureInSermon['bibleVersion'][i] != 'nkjv':
                chapter += ' ' + f"({englishScrpitureInSermon['bibleVersion'][i]})"
            self.setChapter(tf,chapter)
            pageContent = self.seperateVerse(englishVerses[i])
            p = tf.add_paragraph()
            run = p.add_run()
            if len(pageContent) == 1:
                self.represent_verse_in_ppt(pageContent[0],p)
            else:
                for index in range(len(pageContent)):
                    if index == 0:
                        self.represent_verse_in_ppt(pageContent[0],p)
                    else:
                        self.addNewPageToFinishTheRest(pageContent[index])
            font = run.font
            font.size = Pt(54)

            slide = self.prs.slides.add_slide(self.layout)
            tf = self.initTxBox(slide)
            self.setChapter(tf,chineseScrpitureInSermon[i],'chinese')
            pageContent = self.seperateVerse(chineseVerses[i],'chinese')
            p = tf.add_paragraph()
            run = p.add_run()
            if len(pageContent) == 1:
                self.represent_verse_in_ppt(pageContent[0],p,'chinese')
            else:
                for index in range(len(pageContent)):
                    if index == 0:
                        self.represent_verse_in_ppt(pageContent[0], p, 'chinese')
                        content = pageContent[0].split(self.start_verse_token)
                    else:
                        self.addNewPageToFinishTheRest(pageContent[index],'chinese')
            font = run.font
            font.name = '微軟正黑體'
            font.size = Pt(54)

    def initTxBox(self,slide):
        left = top = Inches(0.3)  # 位置
        width = Inches(15)  # 大小
        height = Inches(8)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        return tf

    def setChapter(self,tf,info,language = 'english'):
        p = tf.paragraphs[0]
        book = info.split(' ')[0]
        chapter = info.split(' ')[1].split(':')[0]
        start_verse = info.split(':')[1].split('-')[0]
        end_verse = info.split(':')[1].split('-')[1].split(' ')[0]
        if start_verse == end_verse:
            info = f'{book} {chapter}:{start_verse}'
        run = p.add_run()
        run.text = info
        font = run.font
        font.size = Pt(54)
        font.bold = True
        if language == 'chinese':
            font.name = '微軟正黑體'
        font.color.rgb = RGBColor(255, 0, 0)

    def addPage(self,scripture,verses,language = 'english'):
        slide = self.prs.slides.add_slide(self.layout)
        tf = self.initTxBox(slide)
        self.setChapter(tf,scripture,language)
        p = tf.add_paragraph()

        run = p.add_run()
        pageContent = self.seperateVerse(verses,language)
        if len(pageContent) == 1:
            content = pageContent[0].split(self.start_verse_token)
            for c in content:
                if c != '':
                    scr = c.split(" ")[0]
                    verse = " ".join(c.split(" ")[1:])
                    self.show_content(scr,p,True,language)
                    self.show_content(verse,p,False,language)
        else:
            for index in range(len(pageContent)):
                if index == 0:
                    content = pageContent[index].split(self.start_verse_token)
                    for c in content:
                        if c != '':
                            scr = c.split(" ")[0]
                            verse = " ".join(c.split(" ")[1:])
                            self.show_content(scr, p, True, language)
                            self.show_content(verse, p, False, language)
                else:
                    self.addNewPageToFinishTheRest(pageContent[index],language)
        font = run.font

        if language == 'chinese':
            font.name = '微軟正黑體'
        font.size = Pt(54)


    def addNewPageToFinishTheRest(self,content,language = 'english'):
        slide = self.prs.slides.add_slide(self.layout)
        tf = self.initTxBox(slide)

        p = tf.paragraphs[0]
        content = content.split(self.start_verse_token)
        if content[0] == '':
            for c in content:
                if c != '':
                    scr = c.split(" ")[0]
                    verse = " ".join(c.split(" ")[1:])
                    self.show_content(scr, p, True, language)
                    self.show_content(verse, p, False, language)
        else:
            self.show_content(content[0], p, False, language)
            for index in range(1,len(content)):
                scr = content[index].split(" ")[0]
                verse = " ".join(content[index].split(" ")[1:])
                self.show_content(scr, p, True, language)
                self.show_content(verse, p, False, language)

    def seperateVerse(self,verses,language = 'english'):
        contentEachPage = []
        i = 0
        block = ""
        if language == 'english':
            for index in range(len(verses)):
                words = verses[index].split(" ")
                for word in words:
                    block += word + " "
                    i += 1
                    if i % self.word_limit == 0:
                        contentEachPage.append(block)
                        block = ''
        else:
            for index in range(len(verses)):
                for word in verses[index]:
                    block += word
                    i += 1
                    if i % self.word_limit_chinese == 0:
                        contentEachPage.append(block)
                        block = ''
        contentEachPage.append(block)
        return contentEachPage

    def represent_verse_in_ppt(self,content,p,language = 'english'):
        content = content.split(self.start_verse_token)
        if content[0] == '':
            for c in content:
                if c != '':
                    scr = c.split(" ")[0]
                    verse = " ".join(c.split(" ")[1:])
                    self.show_content(scr, p, True, language)
                    self.show_content(verse, p, False, language)

    def set_superscript(self,font):
        font._element.set('baseline', '30000')

    def show_content(self,content,p,superscript = False,language = 'english'):
        run = p.add_run()
        run.text = content
        font = run.font
        font.size = Pt(54)
        if language == 'chinese':
            font.name = '微軟正黑體'
        if superscript:
            self.set_superscript(font)

    def getEnglishBibleVerses(self,scripture,bibleVersion):
        bibleApi = BibleApi()
        englishBookDic = bibleApi.englishBookDic
        for key in englishBookDic.keys():
            if key in scripture:
                # book = englishBookDic[key]
                book = key
                chapter = scripture.split(" ")[-1].split(":")[0]
                verses = scripture.split(" ")[-1].split(":")[1].split("-")
                startVerse = scripture.split(" ")[-1].split(":")[1].split("-")[0]
                endVerse = startVerse
                if len(verses) != 1:
                    endVerse = scripture.split(" ")[-1].split(":")[1].split("-")[1]
        try:
            return bibleApi.getEnglishVerse(bibleVersion,book,chapter,startVerse,endVerse)
        except:
            pass

    def getChineseBibleVerses(self,info):
        bibleApi = BibleApi()
        chineseBookDic = bibleApi.chineseBookDic
        for key in chineseBookDic.keys():
            if key in info:
                book = chineseBookDic[key]
                chapter = info.split(" ")[-1].split(":")[0]
                verses = info.split(" ")[-1].split(":")[1].split("-")
                startVerse = info.split(" ")[-1].split(":")[1].split("-")[0]
                endVerse = startVerse
                if len(verses) != 1:
                    endVerse = info.split(" ")[-1].split(":")[1].split("-")[1]
        return bibleApi.getChineseVerse(book,chapter,startVerse,endVerse)

    def addClosingSong(self):
        closingSong_slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(closingSong_slide_layout)
        title = slide.shapes.title
        title.text = self.data["closingSongName"]
        title.width = Cm(40)
        title.top = Inches(4)
        title.text_frame.paragraphs[0].font.size = Pt(60)

        try:
            closingSong = self.data["closingSong"]
            content = ""
            contentList = []
            for index in range(len(closingSong["closingSong_chinese"])):
                content += closingSong["closingSong_chinese"][index] + '\n'
                if (index + 1) % 6 == 0:
                    contentList.append(content)
                    content = ""


            contentList.append(content)
            for lyrics in contentList:
                if lyrics != '':
                    self.addNewPageToFinishTheRest(lyrics,'chinese')

            content = ""
            contentList = []
            for index in range(len(closingSong["closingSong_english"])):
                content += closingSong["closingSong_english"][index] + '\n'
                if (index + 1) % 8 == 0:
                    contentList.append(content)
                    content = ""

            contentList.append(content)
            for lyrics in contentList:
                if lyrics != '' and lyrics != '\n':
                    self.addNewPageToFinishTheRest(lyrics)
        except:
            pass


    def addBlessingSong(self):
        blessingSong_slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(blessingSong_slide_layout)
        title = slide.shapes.title
        title.text = "The Blessing Song"
        title.width = Cm(40)
        title.top = Inches(4)
        title.text_frame.paragraphs[0].font.size = Pt(70)

        blessing_song = self.data["blessing_song"]
        contentInFirstPage = []
        contentInSecondPage = []
        for i in range(len(blessing_song)):
            if i <= 6 :
                contentInFirstPage.append(blessing_song[i])
            else:
                contentInSecondPage.append(blessing_song[i])

        slide = self.prs.slides.add_slide(self.layout)
        tf = self.initTxBox(slide)
        p = tf.paragraphs[0]
        run = p.add_run()
        for lyrics in contentInFirstPage:
            run.text += lyrics + '\n'
        font = run.font
        font.size = Pt(60)

        slide = self.prs.slides.add_slide(self.layout)
        tf = self.initTxBox(slide)
        p = tf.paragraphs[0]
        run = p.add_run()
        for lyrics in contentInSecondPage:
            run.text += lyrics + '\n'
        font = run.font
        font.size = Pt(60)


    def addCalvaryImg(self):
        img_path = os.path.dirname(__file__) + '/Img/calvary.png'
        slide = self.prs.slides.add_slide(self.layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)
        top = Inches(0)
        left = Inches(2)
        height = Inches(9)
        slide.shapes.add_picture(img_path, left, top,height=height)

